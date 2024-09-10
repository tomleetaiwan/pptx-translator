import streamlit as st
import os
import openai
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
from dotenv import load_dotenv
from tenacity import (retry,stop_after_attempt,wait_random_exponential,)

# Global variables
total_shapes = 0
shape_count = 0

# Load environment variables from a .env file
load_dotenv()

# Set up the Azure OpenAI Service configuration
openai.api_type = os.getenv("OPENAI_API_TYPE")
openai.api_version = os.getenv("OPENAI_API_VERSION")
API_KEY = os.getenv("OPENAI_API_KEY")
assert API_KEY, "Error: Lack of Azure OpenAI Service API key"
openai.api_key = API_KEY
RESOURCE_ENDPOINT = os.getenv("OPENAI_API_BASE")
assert RESOURCE_ENDPOINT, "Error: Lack of Azure OpenAI service endpoint"
assert "openai.azure.com" in RESOURCE_ENDPOINT.lower(), "Error: The endpoint format of Azure OpenAI Service API is: \n\n\t<Your Azure OpenAI Resource Name>.openai.azure.com"
openai.api_base = RESOURCE_ENDPOINT
model = os.getenv("DEPLOYMENT_NAME")

# Define a retry strategy for the Azure OpenAI Service API call to handle the error of Token per minute limit
@retry(wait=wait_random_exponential(min=60, max=65), stop=stop_after_attempt(10))
def completion_with_backoff(**kwargs):
    return openai.ChatCompletion.create(**kwargs)

# Function to translate text to Chinese using the OpenAI model
def translate_to_chinese(original_text, model, temperature=0.2, max_tokens=2048):
    response = completion_with_backoff(
        engine=model,
        messages=[
            {"role": "user", "content": f'translate "{original_text}" into traditional chinese but the Arabic numerals do not need to be translated, just put translated result without any other descriptions'}
        ],
        temperature=temperature,
        max_tokens=max_tokens
    )
    chinese_text = response["choices"][0]["message"]["content"]
    return chinese_text

# Function to calculate the total number of shapes in a PowerPoint slide
def calculate_shape(shape):
    global total_shapes
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            calculate_shape(sub_shape)
    
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                total_shapes += 1
    else:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    total_shapes += 1
        else:
            total_shapes += 1

# Function to process shapes in a PowerPoint slide
def process_shape(shape, model, progress_bar):    
    global shape_count
    global total_shapes
    
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            process_shape(sub_shape, model,  progress_bar)
    
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                if original_text == "":
                    continue
                # Translate the text to Chinese
                translated_text = translate_to_chinese(original_text, model)
                # Replace the original text with the translated text
                run.text = translated_text
                shape_count += 1
                progress_bar.progress(shape_count / total_shapes)        
    else:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    original_text = cell.text
                    if original_text == "":
                        continue
                    # Translate the text to Chinese
                    translated_text = translate_to_chinese(original_text, model)
                    # Replace the original text with the translated text
                    cell.text = translated_text
                    shape_count += 1
                    progress_bar.progress(shape_count / total_shapes)  
        else:
            shape_count += 1            

            
# Set up the Streamlit app and Azure OpenAI Service configuration
st.set_page_config(page_title="PowerPoint Translator")

# Load the PowerPoint file
uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])
if uploaded_file is not None:
    pr = Presentation(uploaded_file)
    # Count the total number of shapes in the PowerPoint file
    for slide in pr.slides:
        for shape in slide.shapes:
            calculate_shape(shape)    

    # Set up the progress bar
    progress_bar = st.progress(shape_count)
    # Iterate through all the slides and replace the text with Traditional Chinese
    for slide in pr.slides:
        for shape in slide.shapes:
            process_shape(shape, model, progress_bar)

    # Save the translated presentation to a BytesIO object
    buffer = BytesIO()
    pr.save(buffer)
    buffer.seek(0)
    progress_bar.progress(1.0)
    # Download the updated PowerPoint file
    st.download_button(
        label="Download translated file",
        data=buffer,
        file_name="translated.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
